"""
masis.eval.ingest_docs
======================
Document ingestion pipeline for MASIS retrieval.

Ingests documents from the Infosys company docs folder into ChromaDB (vector)
and builds a BM25 index (sparse), then registers both on the researcher module
so the MASIS graph can retrieve evidence.

Pipeline
--------
1. Load documents (LlamaParse v2 for PDF/DOCX/PPTX, direct read for TXT/MD/CSV,
   OpenAI vision for images)
2. Extract document-level metadata via LLM
3. Hierarchical chunking (parent=2000 tokens, child=500 tokens)
4. Table-aware processing (summary for embedding, raw table for context)
5. Embed via text-embedding-3-small and store in ChromaDB
6. Build BM25 index from same corpus
7. Register on masis.agents.researcher

Usage
-----
    python -m masis.eval.ingest_docs                          # default infosys docs
    python -m masis.eval.ingest_docs --doc-folder /path/to/docs --force

"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import json
import logging
import os
import re
import sys
import uuid
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

# ---------------------------------------------------------------------------
# Ensure project root on path for imports
# ---------------------------------------------------------------------------
_PROJECT_ROOT = str(Path(__file__).resolve().parents[1])
if _PROJECT_ROOT not in sys.path:
    sys.path.insert(0, _PROJECT_ROOT)

logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)-8s] %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Settings (from MASIS config or fallback)
# ---------------------------------------------------------------------------
try:
    from masis.config.settings import get_settings
    _settings = get_settings()
    OPENAI_API_KEY = _settings.openai_api_key
    CHROMA_PERSIST_DIR = _settings.chroma_persist_dir
except Exception:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(_PROJECT_ROOT, ".env"))
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
    CHROMA_PERSIST_DIR = os.getenv("CHROMA_PERSIST_DIRECTORY", "./chroma_db")

EMBEDDING_MODEL = os.getenv("MODEL_EMBEDDER", "text-embedding-3-small")
LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY", "")
CHROMA_OPENAI_API_KEY = os.getenv("CHROMA_OPENAI_API_KEY", OPENAI_API_KEY)
if CHROMA_OPENAI_API_KEY:
    os.environ.setdefault("CHROMA_OPENAI_API_KEY", CHROMA_OPENAI_API_KEY)


def _resolve_openai_api_key() -> str:
    """Resolve API key for Chroma OpenAI embeddings with robust fallback order."""
    key = (
        os.getenv("CHROMA_OPENAI_API_KEY")
        or os.getenv("OPENAI_API_KEY")
        or CHROMA_OPENAI_API_KEY
        or OPENAI_API_KEY
    )
    if key:
        os.environ.setdefault("OPENAI_API_KEY", key)
        os.environ.setdefault("CHROMA_OPENAI_API_KEY", key)
        return key

    try:
        from dotenv import load_dotenv

        load_dotenv(os.path.join(_PROJECT_ROOT, ".env"))
    except Exception:
        pass

    key = os.getenv("CHROMA_OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY") or ""
    if key:
        os.environ.setdefault("OPENAI_API_KEY", key)
        os.environ.setdefault("CHROMA_OPENAI_API_KEY", key)
    return key


def _normalize_persist_dir(persist_dir: str) -> str:
    """Normalize Chroma path to avoid duplicate-client identity mismatches."""
    return os.path.abspath(os.path.normpath(persist_dir))


def _get_or_create_chroma_client(persist_dir: str) -> Any:
    """Return one persistent Chroma client per normalized path."""
    import chromadb
    from chromadb.api.shared_system_client import SharedSystemClient

    global _CHROMA_CLIENT, _CHROMA_CLIENT_PATH
    normalized = _normalize_persist_dir(persist_dir)
    if _CHROMA_CLIENT is not None and _CHROMA_CLIENT_PATH == normalized:
        return _CHROMA_CLIENT

    os.makedirs(normalized, exist_ok=True)
    try:
        _CHROMA_CLIENT = chromadb.PersistentClient(path=normalized)
    except ValueError as exc:
        # Streamlit reruns (or mixed client settings in-process) can leave a
        # cached SharedSystemClient for the same path with different settings.
        # Recover by clearing Chroma's process-local system cache and retrying.
        msg = str(exc)
        if "already exists" in msg and "different settings" in msg:
            logger.warning(
                "Chroma client settings collision detected for %s; clearing "
                "SharedSystemClient cache and retrying once.",
                normalized,
            )
            SharedSystemClient.clear_system_cache()
            _CHROMA_CLIENT = chromadb.PersistentClient(path=normalized)
        else:
            raise
    _CHROMA_CLIENT_PATH = normalized
    return _CHROMA_CLIENT

# Chunking parameters
CHUNK_SIZE_PARENT = 2000
CHUNK_SIZE_CHILD = 500
CHUNK_OVERLAP = 50
LLAMAPARSE_TIER = "agentic"

def _resolve_default_doc_folder() -> str:
    """Resolve a default Infosys document folder from common local layouts."""
    project_root = Path(__file__).resolve().parents[2]
    candidates = [project_root / "infosys_company_docs"]
    candidates.extend(
        project_root.glob("rag_pipeline*/Insights-Agent-Flow-Research-main/infosys_company_docs")
    )
    for path in candidates:
        if path.exists():
            return str(path)
    return str(project_root / "infosys_company_docs")


# Default document folder
DEFAULT_DOC_FOLDER = _resolve_default_doc_folder()

CHROMA_COLLECTION_NAME = "masis_eval_infosys"
PARENT_STORE_FILENAME = "parent_store_infosys.json"

# Module-level ChromaDB client — must stay alive for registered collection to work
_CHROMA_CLIENT = None
_CHROMA_CLIENT_PATH = ""

# File extension sets
TEXT_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}
DATA_EXTS = {".csv"}
ALL_SUPPORTED_EXTS = TEXT_EXTS | IMAGE_EXTS | DATA_EXTS


# ===========================================================================
# Step 0: Hash-based caching
# ===========================================================================

def _compute_file_hash(filepath: str) -> str:
    """Compute SHA-256 hash of a file for change detection."""
    sha256 = hashlib.sha256()
    with open(filepath, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            sha256.update(block)
    return sha256.hexdigest()


def _compute_folder_hashes(doc_folder: str) -> Dict[str, str]:
    """Compute hashes for all supported files in a document folder."""
    hashes: Dict[str, str] = {}
    for root, _, files in os.walk(doc_folder):
        for name in sorted(files):
            ext = os.path.splitext(name)[1].lower()
            if ext in ALL_SUPPORTED_EXTS:
                filepath = os.path.join(root, name)
                rel_path = os.path.relpath(filepath, doc_folder)
                hashes[rel_path] = _compute_file_hash(filepath)
    return hashes


# ===========================================================================
# Step 1: Load documents  --  multi-format parsing
# ===========================================================================

class _Document:
    """Lightweight document wrapper (avoids llama_index dependency)."""

    def __init__(self, text: str, metadata: Optional[Dict[str, Any]] = None):
        self.text = text
        self.metadata = metadata or {}


def _load_pdf(filepath: str) -> List[_Document]:
    """Load PDF via pypdf (no external API needed)."""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed  --  trying LlamaParse for %s", filepath)
        return _load_via_llamaparse_sync(filepath)

    try:
        reader = PdfReader(filepath)
        docs: List[_Document] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                docs.append(_Document(
                    text=text,
                    metadata={
                        "file_name": os.path.basename(filepath),
                        "file_path": filepath,
                        "page_label": i,
                        "source": "pypdf",
                    },
                ))
        logger.info("  [pypdf] %s  ->  %d page(s)", os.path.basename(filepath), len(docs))
        return docs
    except Exception as exc:
        logger.error("  [pypdf] Failed to read %s: %s", filepath, exc)
        return _load_via_llamaparse_sync(filepath)


def _load_via_llamaparse_sync(filepath: str) -> List[_Document]:
    """Fallback: use LlamaParse v2 if available and API key set."""
    if not LLAMA_CLOUD_API_KEY:
        logger.warning("  No LLAMA_CLOUD_API_KEY  --  skipping %s", filepath)
        return []

    try:
        from llama_cloud import AsyncLlamaCloud
    except ImportError:
        logger.warning("  llama_cloud not installed  --  skipping %s", filepath)
        return []

    async def _parse() -> List[_Document]:
        client = AsyncLlamaCloud(api_key=LLAMA_CLOUD_API_KEY)
        filename = os.path.basename(filepath)
        with open(filepath, "rb") as fh:
            file_obj = await client.files.create(file=(filename, fh), purpose="parse")
        result = await client.parsing.parse(
            file_id=file_obj.id,
            tier=LLAMAPARSE_TIER,
            version="latest",
            output_options={
                "markdown": {"tables": {"output_tables_as_markdown": True}},
            },
            expand=["markdown"],
        )
        docs: List[_Document] = []
        for page in result.markdown.pages:
            page_text = page.markdown or ""
            if page_text.strip():
                docs.append(_Document(
                    text=page_text,
                    metadata={
                        "file_name": filename,
                        "file_path": filepath,
                        "page_label": page.page_number,
                        "source": "llamaparse_v2",
                    },
                ))
        return docs

    try:
        return asyncio.run(_parse())
    except RuntimeError:
        import concurrent.futures
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(asyncio.run, _parse())
            return future.result()


def _load_docx(filepath: str) -> List[_Document]:
    """Load DOCX via python-docx."""
    try:
        from docx import Document as DocxDoc
    except ImportError:
        logger.warning("python-docx not installed  --  skipping %s", filepath)
        return _load_via_llamaparse_sync(filepath)

    try:
        doc = DocxDoc(filepath)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if not text.strip():
            return []
        logger.info("  [docx] %s  ->  %d chars", os.path.basename(filepath), len(text))
        return [_Document(
            text=text,
            metadata={
                "file_name": os.path.basename(filepath),
                "file_path": filepath,
                "page_label": 1,
                "source": "python_docx",
            },
        )]
    except Exception as exc:
        logger.error("  [docx] Failed to read %s: %s", filepath, exc)
        return []


def _load_pptx(filepath: str) -> List[_Document]:
    """Load PPTX via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed  --  skipping %s", filepath)
        return _load_via_llamaparse_sync(filepath)

    try:
        prs = Presentation(filepath)
        docs: List[_Document] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                docs.append(_Document(
                    text=slide_text,
                    metadata={
                        "file_name": os.path.basename(filepath),
                        "file_path": filepath,
                        "page_label": i,
                        "source": "python_pptx",
                    },
                ))
        logger.info("  [pptx] %s  ->  %d slide(s)", os.path.basename(filepath), len(docs))
        return docs
    except Exception as exc:
        logger.error("  [pptx] Failed to read %s: %s", filepath, exc)
        return []


def _load_text_file(filepath: str) -> List[_Document]:
    """Load plain text or markdown files."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        if not text.strip():
            return []
        return [_Document(
            text=text,
            metadata={
                "file_name": os.path.basename(filepath),
                "file_path": filepath,
                "page_label": 1,
                "source": "simple_reader",
            },
        )]
    except Exception as exc:
        logger.error("Failed to read text file %s: %s", filepath, exc)
        return []


def _load_csv(filepath: str) -> List[_Document]:
    """Load CSV as text (with header preservation)."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        if not text.strip():
            return []
        return [_Document(
            text=text,
            metadata={
                "file_name": os.path.basename(filepath),
                "file_path": filepath,
                "page_label": 1,
                "source": "csv_reader",
            },
        )]
    except Exception as exc:
        logger.error("Failed to read CSV %s: %s", filepath, exc)
        return []


def _caption_image(filepath: str) -> str:
    """Generate a semantic caption for an image using OpenAI vision."""
    if not OPENAI_API_KEY:
        logger.warning("No OPENAI_API_KEY  --  skipping image captioning for %s", filepath)
        return ""

    try:
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        with open(filepath, "rb") as f:
            image_bytes = f.read()
        data_url = "data:image/png;base64," + base64.b64encode(image_bytes).decode("utf-8")

        response = client.chat.completions.create(
            model="gpt-4.1",
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Analyze this image in detail. If it's a chart or graph, "
                            "extract all data points, labels, and trends. If it's a table, "
                            "convert it to markdown format. If it's a diagram, describe the "
                            "structure and relationships. Provide a comprehensive textual "
                            "description that captures all information in the image."
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }],
        )
        caption = response.choices[0].message.content.strip() if response.choices[0].message.content else ""
        logger.info("  [vision] %s  ->  %d chars caption", os.path.basename(filepath), len(caption))
        return caption
    except Exception as exc:
        logger.warning("  [vision] Captioning failed for %s: %s", filepath, exc)
        return ""


def _load_image(filepath: str) -> List[_Document]:
    """Load an image file by captioning it with OpenAI vision."""
    text = _caption_image(filepath)
    if not text.strip():
        return []
    return [_Document(
        text=text,
        metadata={
            "file_name": os.path.basename(filepath),
            "file_path": filepath,
            "page_label": 0,
            "source": "vision_caption",
        },
    )]


def load_documents(doc_folder: str) -> List[_Document]:
    """
    Load all documents from doc_folder.

    Dispatches by extension:
        .pdf          ->  pypdf (fallback: LlamaParse v2)
        .docx         ->  python-docx (fallback: LlamaParse v2)
        .pptx         ->  python-pptx (fallback: LlamaParse v2)
        .txt / .md    ->  direct read
        .csv          ->  direct read
        .png / .jpg   ->  OpenAI vision captioning
    """
    if not os.path.exists(doc_folder):
        raise FileNotFoundError(f"Document folder not found: {doc_folder}")

    all_docs: List[_Document] = []
    file_count = 0

    for root, _, files in os.walk(doc_folder):
        for name in sorted(files):
            ext = os.path.splitext(name)[1].lower()
            if ext not in ALL_SUPPORTED_EXTS:
                continue

            filepath = os.path.join(root, name)
            file_count += 1
            logger.info("Loading [%d]: %s", file_count, name)

            if ext == ".pdf":
                all_docs.extend(_load_pdf(filepath))
            elif ext == ".docx":
                all_docs.extend(_load_docx(filepath))
            elif ext == ".pptx":
                all_docs.extend(_load_pptx(filepath))
            elif ext in {".txt", ".md"}:
                all_docs.extend(_load_text_file(filepath))
            elif ext == ".csv":
                all_docs.extend(_load_csv(filepath))
            elif ext in IMAGE_EXTS:
                all_docs.extend(_load_image(filepath))

    logger.info("Total: %d document pages from %d files", len(all_docs), file_count)
    return all_docs


# ===========================================================================
# Step 2: Metadata extraction
# ===========================================================================

def _extract_document_metadata(text_first_pages: str) -> Dict[str, Any]:
    """Extract document-level metadata (type, year, quarter) via LLM."""
    try:
        from langchain_openai import ChatOpenAI
    except ImportError:
        logger.warning("langchain_openai not installed  --  using keyword-only metadata")
        return _extract_metadata_keywords(text_first_pages)

    if not OPENAI_API_KEY:
        return _extract_metadata_keywords(text_first_pages)

    try:
        llm = ChatOpenAI(model="gpt-4.1-mini", temperature=0)
        prompt = (
            "You are a document metadata extraction assistant.\n"
            "Infer normalized metadata from the excerpt and return strict JSON only.\n\n"
            "Required keys:\n"
            '- "document_type": one of "annual_report", "quarterly_report", '
            '"strategy_note", "operational_update", "press_release", "other"\n'
            '- "year": 4-digit year string or ""\n'
            '- "quarter": one of "Q1","Q2","Q3","Q4" or ""\n\n'
            "Extraction rules:\n"
            "- Prefer explicit statements over inference.\n"
            "- If uncertain, return empty string for year/quarter.\n"
            "- Keep output machine-parseable; no markdown, no prose.\n\n"
            f"Document excerpt:\n{text_first_pages[:3000]}\n\n"
            "Return ONLY valid JSON."
        )
        response = llm.invoke(prompt)
        content = response.content.strip()
        json_match = re.search(r'```(?:json)?\s*(.*?)\s*```', content, re.DOTALL)
        if json_match:
            content = json_match.group(1)
        return json.loads(content)
    except Exception as exc:
        logger.warning("LLM metadata extraction failed: %s  --  using keywords", exc)
        return _extract_metadata_keywords(text_first_pages)


def _extract_metadata_keywords(text: str) -> Dict[str, Any]:
    """Keyword-based metadata extraction (fallback)."""
    meta: Dict[str, Any] = {"document_type": "other", "year": "", "quarter": ""}
    text_lower = text.lower()

    year_match = re.search(r"\b(20\d{2})\b", text)
    if year_match:
        meta["year"] = year_match.group(1)

    quarter_match = re.search(r"\b(q[1-4])\b", text_lower)
    if quarter_match:
        meta["quarter"] = quarter_match.group(1).upper()

    if any(kw in text_lower for kw in ["press release", "ifrs", "results for"]):
        meta["document_type"] = "press_release"
    elif any(kw in text_lower for kw in ["quarterly", "quarter ended"]):
        meta["document_type"] = "quarterly_report"
    elif "annual" in text_lower:
        meta["document_type"] = "annual_report"

    return meta


def _detect_section(text: str) -> str:
    """Detect section from chunk text."""
    section_keywords = {
        "Financial Performance": ["revenue", "income", "profit", "loss", "earnings", "margin"],
        "Risk Factors": ["risk", "uncertainty", "challenge", "threat", "concern"],
        "Operations": ["operations", "operational", "production", "efficiency"],
        "Strategy": ["strategy", "strategic", "initiative", "roadmap", "vision"],
        "Personnel": ["employee", "headcount", "hiring", "attrition", "talent", "workforce"],
        "Technology": ["technology", "digital", "platform", "AI", "generative"],
        "Compliance": ["compliance", "regulatory", "regulation", "audit", "governance"],
        "Market": ["market", "competition", "competitor", "industry"],
    }
    text_lower = text.lower()
    scores = {}
    for section, keywords in section_keywords.items():
        score = sum(1 for kw in keywords if kw in text_lower)
        if score > 0:
            scores[section] = score
    return max(scores, key=scores.get) if scores else "General"


def _detect_department(text: str) -> str:
    """Detect department from chunk text."""
    dept_keywords = {
        "Sales": ["sales", "selling", "revenue generation", "bookings"],
        "Engineering": ["engineering", "development", "r&d", "software"],
        "Marketing": ["marketing", "brand", "campaign", "advertising"],
        "Finance": ["finance", "accounting", "budget", "treasury"],
        "HR": ["human resources", "hr", "talent", "recruitment", "employee"],
        "Operations": ["operations", "supply chain", "logistics"],
        "Legal": ["legal", "compliance", "regulatory", "contract"],
        "IT": ["information technology", "cybersecurity", "data center"],
    }
    text_lower = text.lower()
    scores = {}
    for dept, keywords in dept_keywords.items():
        score = sum(1 for kw in keywords if kw in text_lower)
        if score > 0:
            scores[dept] = score
    return max(scores, key=scores.get) if scores else "Company-wide"


# ===========================================================================
# Step 3: Hierarchical chunking
# ===========================================================================

def _text_split(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Simple sentence-aware text splitter."""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []

    # Try to split on sentence boundaries
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks: List[str] = []
    current = ""

    for sentence in sentences:
        if len(current) + len(sentence) + 1 > chunk_size and current:
            chunks.append(current.strip())
            # Keep overlap from end of current chunk
            current = current[-overlap:] + " " + sentence if overlap > 0 else sentence
        else:
            current = current + " " + sentence if current else sentence

    if current.strip():
        chunks.append(current.strip())

    return chunks


def _is_table_content(text: str) -> bool:
    """Detect markdown tables."""
    lines = text.strip().split("\n")
    pipe_lines = [line for line in lines if "|" in line]
    separator_lines = [line for line in lines if re.match(r'^[\s|:-]+$', line)]
    return len(pipe_lines) >= 3 and len(separator_lines) >= 1


def _extract_tables(text: str) -> List[Dict[str, Any]]:
    """Extract table blocks from text."""
    table_pattern = re.compile(
        r'((?:\|[^\n]+\|\n)+(?:\|[-:\s|]+\|\n)(?:\|[^\n]+\|\n)*)',
        re.MULTILINE,
    )
    tables: List[Dict[str, Any]] = []
    last_end = 0
    for match in table_pattern.finditer(text):
        start, end = match.span()
        tables.append({
            "table": match.group(0).strip(),
            "before": text[last_end:start].strip(),
            "start": start,
            "end": end,
        })
        last_end = end
    return tables


def create_hierarchical_chunks(
    text: str,
    source_file: str,
    page_num: int = 0,
    doc_metadata: Optional[Dict[str, Any]] = None,
) -> List[Dict[str, Any]]:
    """
    Create parent and child chunks with linking.

    Child chunks (~500 chars) are embedded in the vector store.
    Parent chunks (~2000 chars) are stored separately and returned as LLM context.
    """
    if not text.strip():
        return []

    doc_metadata = doc_metadata or {}
    chunks: List[Dict[str, Any]] = []

    # Check for tables first
    tables = _extract_tables(text)

    if tables:
        last_end = 0
        for table_info in tables:
            # Text before this table
            text_before = text[last_end:table_info["start"]].strip()
            if text_before:
                chunks.extend(
                    _chunk_text_hierarchical(text_before, source_file, page_num, doc_metadata)
                )

            # Table chunk: use table text as both embedded and context
            table_id = str(uuid.uuid4())
            parent_id = str(uuid.uuid4())
            chunk = {
                "id": table_id,
                "text": table_info["table"],
                "parent_id": parent_id,
                "parent_text": table_info["table"],
                "content_type": "table",
                "source_file": source_file,
                "page_number": page_num,
                "document_type": doc_metadata.get("document_type", "other"),
                "year": str(doc_metadata.get("year", "")),
                "quarter": str(doc_metadata.get("quarter", "")),
                "section": _detect_section(table_info["table"]),
                "department": _detect_department(table_info["table"]),
            }
            chunks.append(chunk)
            last_end = table_info["end"]

        # Text after last table
        remaining = text[last_end:].strip()
        if remaining:
            chunks.extend(
                _chunk_text_hierarchical(remaining, source_file, page_num, doc_metadata)
            )
    else:
        chunks.extend(
            _chunk_text_hierarchical(text, source_file, page_num, doc_metadata)
        )

    return chunks


def _chunk_text_hierarchical(
    text: str,
    source_file: str,
    page_num: int,
    doc_metadata: Dict[str, Any],
) -> List[Dict[str, Any]]:
    """Create parent-child chunk pairs from plain text."""
    parent_texts = _text_split(text, CHUNK_SIZE_PARENT, CHUNK_OVERLAP * 2)
    chunks: List[Dict[str, Any]] = []

    for parent_text in parent_texts:
        parent_id = str(uuid.uuid4())
        child_texts = _text_split(parent_text, CHUNK_SIZE_CHILD, CHUNK_OVERLAP)

        for child_text in child_texts:
            chunk = {
                "id": str(uuid.uuid4()),
                "text": child_text,
                "parent_id": parent_id,
                "parent_text": parent_text,
                "content_type": "text",
                "source_file": source_file,
                "page_number": page_num,
                "document_type": doc_metadata.get("document_type", "other"),
                "year": str(doc_metadata.get("year", "")),
                "quarter": str(doc_metadata.get("quarter", "")),
                "section": _detect_section(child_text),
                "department": _detect_department(child_text),
            }
            chunks.append(chunk)

    return chunks


# ===========================================================================
# Step 4: Store in ChromaDB + Build BM25
# ===========================================================================

def store_in_chromadb(
    chunks: List[Dict[str, Any]],
    persist_dir: str,
    collection_name: str,
) -> Tuple[Any, Dict[str, int]]:
    """
    Store chunks in ChromaDB with OpenAI embeddings.

    Returns (collection, stats_dict).
    """
    import chromadb
    from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction

    if not chunks:
        logger.warning("No chunks to store!")
        return None, {"num_chunks": 0, "num_parents": 0}

    key = _resolve_openai_api_key()
    if not key:
        raise RuntimeError(
            "OpenAI API key is missing. Set OPENAI_API_KEY (or CHROMA_OPENAI_API_KEY) "
            "in environment or .env before running retrieval setup."
        )

    normalized_persist_dir = _normalize_persist_dir(persist_dir)
    client = _get_or_create_chroma_client(normalized_persist_dir)

    # Drop existing collection for fresh ingest
    try:
        client.delete_collection(collection_name)
        logger.info("Deleted existing collection: %s", collection_name)
    except Exception:
        pass

    embedding_fn = OpenAIEmbeddingFunction(
        api_key=key,
        model_name=EMBEDDING_MODEL,
    )

    collection = client.get_or_create_collection(
        name=collection_name,
        embedding_function=embedding_fn,
    )

    parent_store: Dict[str, str] = {}
    ids: List[str] = []
    documents: List[str] = []
    metadatas: List[Dict[str, Any]] = []

    for chunk in chunks:
        ids.append(chunk["id"])
        documents.append(chunk["text"])
        parent_store[chunk["parent_id"]] = chunk["parent_text"]

        metadatas.append({
            "chunk_id": chunk["id"],
            "doc_id": chunk.get("source_file", "unknown"),
            "parent_chunk_id": chunk["parent_id"],
            "content_type": chunk.get("content_type", "text"),
            "source_file": chunk.get("source_file", "unknown"),
            "source_label": chunk.get("source_file", ""),
            "page_number": chunk.get("page_number", 0),
            "document_type": chunk.get("document_type", "other"),
            "year": str(chunk.get("year", "")),
            "quarter": str(chunk.get("quarter", "")),
            "section": chunk.get("section", "General"),
            "department": chunk.get("department", "Company-wide"),
        })

    # Add in batches
    batch_size = 100
    for i in range(0, len(ids), batch_size):
        end = min(i + batch_size, len(ids))
        collection.add(
            ids=ids[i:end],
            documents=documents[i:end],
            metadatas=metadatas[i:end],
        )
        logger.info("  Embedded batch %d: %d chunks", i // batch_size + 1, end - i)

    # Persist parent store
    parent_path = os.path.join(normalized_persist_dir, PARENT_STORE_FILENAME)
    with open(parent_path, "w", encoding="utf-8") as f:
        json.dump(parent_store, f)

    stats = {"num_chunks": len(ids), "num_parents": len(parent_store)}
    logger.info("Stored %d chunks, %d parent chunks", stats["num_chunks"], stats["num_parents"])
    return collection, stats


def build_bm25_index(
    chunks: List[Dict[str, Any]],
) -> Optional[Tuple[Any, List[str], List[Dict[str, Any]]]]:
    """Build a BM25 index from the chunk corpus."""
    try:
        from rank_bm25 import BM25Okapi
    except ImportError:
        logger.warning("rank_bm25 not installed  --  BM25 index not built")
        return None

    corpus = [chunk["text"] for chunk in chunks]
    corpus_meta = [
        {
            "chunk_id": chunk["id"],
            "doc_id": chunk.get("source_file", "unknown"),
            "parent_chunk_id": chunk.get("parent_id"),
            "source_label": chunk.get("source_file", ""),
        }
        for chunk in chunks
    ]

    tokenized_corpus = [doc.lower().split() for doc in corpus]
    bm25_index = BM25Okapi(tokenized_corpus)
    logger.info("Built BM25 index: %d documents", len(corpus))
    return bm25_index, corpus, corpus_meta


# ===========================================================================
# Step 5: Register on researcher module
# ===========================================================================

def register_on_researcher(collection: Any, bm25_data: Optional[Tuple]) -> None:
    """
    Register the ChromaDB collection and BM25 index on the researcher module
    so that the MASIS graph can retrieve evidence.
    """
    try:
        from masis.agents.researcher import set_chroma_collection, set_bm25_index
    except ImportError:
        logger.warning("Could not import researcher module  --  skipping registration")
        return

    set_chroma_collection(collection)
    logger.info("Registered ChromaDB collection on researcher module")

    if bm25_data is not None:
        bm25_index, corpus, corpus_meta = bm25_data
        set_bm25_index(bm25_index, corpus, corpus_meta)
        logger.info("Registered BM25 index on researcher module")


# ===========================================================================
# Step 6: Full pipeline
# ===========================================================================

def setup_retrieval(
    doc_folder: Optional[str] = None,
    persist_dir: Optional[str] = None,
    force: bool = False,
) -> Dict[str, Any]:
    """
    Full ingestion pipeline: load  ->  chunk  ->  embed  ->  register.

    This is the main entry point for setting up retrieval before running
    MASIS queries. Called by scenario_tests.py and regression.py.

    Args:
        doc_folder: Path to document folder. Defaults to infosys_company_docs.
        persist_dir: ChromaDB persistence directory. Defaults to ./chroma_db.
        force: If True, re-ingest even if no changes detected.

    Returns:
        Dict with ingestion stats.
    """
    doc_folder = doc_folder or DEFAULT_DOC_FOLDER
    persist_dir = _normalize_persist_dir(persist_dir or CHROMA_PERSIST_DIR)

    logger.info("=" * 70)
    logger.info("MASIS Document Ingestion Pipeline")
    logger.info("=" * 70)
    logger.info("  Document folder: %s", doc_folder)
    logger.info("  ChromaDB dir:    %s", persist_dir)

    # Check if we can skip ingestion (already ingested, no changes)
    global _CHROMA_CLIENT
    if not force:
        try:
            _CHROMA_CLIENT = _get_or_create_chroma_client(persist_dir)
            existing = _CHROMA_CLIENT.get_collection(
                name=CHROMA_COLLECTION_NAME,
                embedding_function=None,
            )
            count = existing.count()
            if count > 0:
                logger.info("Found existing ChromaDB collection with %d chunks  --  reusing", count)
                # Still need to build BM25 and register
                all_data = existing.get(include=["documents", "metadatas"])
                chunks_for_bm25 = []
                for i, doc_text in enumerate(all_data["documents"]):
                    meta = all_data["metadatas"][i] if all_data["metadatas"] else {}
                    chunks_for_bm25.append({
                        "id": all_data["ids"][i],
                        "text": doc_text,
                        "source_file": meta.get("source_file", "unknown"),
                        "parent_id": meta.get("parent_chunk_id", ""),
                    })

                # Re-create collection with embedding function for queries
                key = _resolve_openai_api_key()
                if not key:
                    raise RuntimeError(
                        "OpenAI API key is missing. Set OPENAI_API_KEY (or CHROMA_OPENAI_API_KEY) "
                        "in environment or .env before running retrieval setup."
                    )
                from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
                embedding_fn = OpenAIEmbeddingFunction(
                    api_key=key,
                    model_name=EMBEDDING_MODEL,
                )
                collection = _CHROMA_CLIENT.get_collection(
                    name=CHROMA_COLLECTION_NAME,
                    embedding_function=embedding_fn,
                )

                bm25_data = build_bm25_index(chunks_for_bm25)
                register_on_researcher(collection, bm25_data)
                return {
                    "num_chunks": count,
                    "skipped": True,
                    "reason": "Existing collection reused",
                }
        except Exception:
            pass  # Collection doesn't exist, proceed with ingestion

    # Step 1: Load documents
    documents = load_documents(doc_folder)
    if not documents:
        return {"num_documents": 0, "num_chunks": 0, "error": "No documents found"}

    # Step 2: Group by file and extract metadata
    docs_by_file: Dict[str, List[_Document]] = {}
    for doc in documents:
        file_name = doc.metadata.get("file_name", "unknown")
        docs_by_file.setdefault(file_name, []).append(doc)

    all_chunks: List[Dict[str, Any]] = []

    for file_name, file_docs in docs_by_file.items():
        logger.info("Processing: %s (%d page(s))", file_name, len(file_docs))

        # Extract doc-level metadata from first few pages
        first_pages_text = "\n\n".join(doc.text for doc in file_docs[:3])[:5000]
        doc_metadata = _extract_document_metadata(first_pages_text)
        logger.info(
            "  Metadata: type=%s, year=%s, quarter=%s",
            doc_metadata.get("document_type"),
            doc_metadata.get("year"),
            doc_metadata.get("quarter"),
        )

        # Step 3: Chunk each page
        for doc in file_docs:
            page_chunks = create_hierarchical_chunks(
                doc.text,
                source_file=file_name,
                page_num=doc.metadata.get("page_label", 0),
                doc_metadata=doc_metadata,
            )
            all_chunks.extend(page_chunks)

    logger.info("Created %d total chunks from %d document(s)", len(all_chunks), len(documents))

    # Step 4: Embed and store
    collection, stats = store_in_chromadb(all_chunks, persist_dir, CHROMA_COLLECTION_NAME)

    # Step 5: Build BM25
    bm25_data = build_bm25_index(all_chunks)

    # Step 6: Register on researcher
    if collection is not None:
        register_on_researcher(collection, bm25_data)

    logger.info("=" * 70)
    logger.info("Ingestion Complete: %d files  ->  %d chunks", len(docs_by_file), len(all_chunks))
    logger.info("=" * 70)

    return {
        "num_documents": len(documents),
        "num_files": len(docs_by_file),
        "skipped": False,
        **stats,
    }


# ===========================================================================
# CLI entry point
# ===========================================================================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="MASIS document ingestion pipeline")
    parser.add_argument(
        "--doc-folder",
        default=DEFAULT_DOC_FOLDER,
        help="Path to document folder (default: infosys_company_docs)",
    )
    parser.add_argument("--persist-dir", default=CHROMA_PERSIST_DIR)
    parser.add_argument("--force", action="store_true", help="Force re-ingestion")
    args = parser.parse_args()

    result = setup_retrieval(
        doc_folder=args.doc_folder,
        persist_dir=args.persist_dir,
        force=args.force,
    )
    print(json.dumps(result, indent=2))
